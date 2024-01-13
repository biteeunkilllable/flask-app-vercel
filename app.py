from flask import Flask, request, jsonify
from pydantic import BaseModel, Field
import PyPDF2
from pptx import Presentation
import io
from docx import Document
from flask_cors import CORS
from dotenv import load_dotenv
import os
import google.generativeai as genai

load_dotenv()
app = Flask(__name__)
CORS(app)

@app.route("/",methods=["GET"])
def defaullt():
    return "Hola Mundo"


genai.configure(api_key=os.environ.get("GOOGLE_API_KEY"))
model = genai.GenerativeModel("gemini-pro")

@app.route('/note', methods=['POST'])
def nt():
    if 'text' in request.form and 'lang' in request.form:
        text = request.form['text']
        lang = request.form['lang']
        text = text.lower()
        text = text.replace("\n\n", "\n")
        result = model.generate_content(
            f"Create 5 notes using dashed points for the following text in {lang}:\n {text}",
            generation_config={
                "max_output_tokens": 8192,
                "temperature": 0.3,
            },)
        token = len(text) + len(result.text)
        print(token)
        return result.text.replace("\n","<br>")
    else:
        return 'Missing required parameters'


@app.route('/translate', methods=['POST'])
def translate_text():
    if 'text' in request.form and 'lang' in request.form:
        text = request.form['text']
        lang = request.form['lang']
        text = text.replace("\n\n", "\n")
        result = model.generate_content(
        f"Translate this text into {lang} :\n {text}",
        generation_config={
            "max_output_tokens": 8192,
            "temperature": 0.0,
        },)
        token = len(text) + len(result.text)
        print(token)
        return result.text.replace("\n","<br>")
    else:
        return 'Missing required parameters'


@app.route('/summarize', methods=['POST'])
def summarize_text():
    if 'text' in request.form and 'lang' in request.form:
        text = request.form['text']
        lang = request.form['lang']
        text = text.replace("\n\n", "\n")
        result = model.generate_content(
            f"Provide a 2 sentence summary of the following text in {lang}:\n {text}",
            generation_config={
                "max_output_tokens": 8192,
                "temperature": 0.2,
            },)
        token = len(text) + len(result.text)
        print(token)
        return result.text.replace("\n","<br>")
    else:
        return 'Missing required parameter'


@app.route('/explain', methods=['POST'])
def explain_text():
    if 'text' in request.form:
        text = request.form['text']
        lang = request.form['lang']
        text = text.replace("\n\n", "\n")
        result = model.generate_content(
        f"Provide a 4 sentence explanation for the following text in {lang}:\n {text}",
        generation_config={
            "max_output_tokens": 8192,
            "temperature": 0.2,
        },)
        token = len(text) + len(result.text)
        print(token)
        return result.text.replace("\n","<br>")
    else:
        return 'Missing required parameter: text'

@app.route('/flash_cards', methods=['POST'])
def FL():

    file = request.files['file']

    if file.filename.split('.')[1]=='txt':
        text = file.read()
        text = text.decode('utf-8')

    elif file.filename.split('.')[1]=='pdf':
        text = ''
        pdf_reader = PyPDF2.PdfReader(file)
        for p in range(2, len(pdf_reader.pages)):
            page = pdf_reader.pages[p]
            text += page.extract_text()

    elif file.filename.split('.')[1] == 'docx':
        docx_content = file.read()
        document = Document(io.BytesIO(docx_content))
        text = ''
        for paragraph in document.paragraphs:
            text+= paragraph.text + '\n'

    elif file.filename.split('.')[1] == 'pptx':
        pptx_content = file.read()
        presentation = Presentation(io.BytesIO(pptx_content))
        text = ''
        p_1 = True
        p_2 = True
        for slide in presentation.slides:
            if p_1:
                p_1 = False
                continue
            elif p_2:
                p_2 = False
                continue
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text+= shape.text+'\n'

    else:
        return 'unsupported file'

    text = text.lower().replace("\n\n\n", "\n").replace("\n\n", "\n")

    fl_num = 15
    l = len(text.split(" "))
    print(l)
    if l > 5000 and l <= 10000:
        fl_num = 25
    elif l > 10000:
        return "Please choose a smaller file for upload."

    format_i = '''
    Follow this format:

    [
        { "question": "here is the question",
            "answer": "here is the answer"
        },
        { "question": "here is the question",
            "answer": "here is the answer"
        }
    ]
    '''
    prompt = f"Create flashcards with {fl_num} question with short answers without duplicate for the following text:\n{text}\n{format_i}"

    result = model.generate_content(
    prompt,
    generation_config={
        "max_output_tokens": 8192,
        "temperature": 0.2,
        },
    )
    token = len(text) + len(result.text)
    print(token)
    return eval(result.text)
    

@app.route('/quizz', methods=['POST'])
def Quizz():
    file = request.files['file']
    num_qus = request.form['num_qus']
    level = request.form['level']

    if file.filename.split('.')[1] == 'txt':
        text = file.read()
        text = text.decode('utf-8')

    elif file.filename.split('.')[1] == 'pdf':
        text = ''
        pdf_reader = PyPDF2.PdfReader(file)
        for p in range(2, len(pdf_reader.pages)):
            page = pdf_reader.pages[p]
            text += page.extract_text()

    elif file.filename.split('.')[1] == 'docx':
        docx_content = file.read()
        document = Document(io.BytesIO(docx_content))
        text = ''
        for paragraph in document.paragraphs:
            text += paragraph.text + '\n'

    elif file.filename.split('.')[1] == 'pptx':
        pptx_content = file.read()
        presentation = Presentation(io.BytesIO(pptx_content))
        text = ''
        p_1 = True
        p_2 = True
        for slide in presentation.slides:
            if p_1:
                p_1 = False
                continue
            elif p_2:
                p_2 = False
                continue
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + '\n'

    else:
        return 'unsupported file'

    text = text.lower().replace("\n\n\n", "\n").replace("\n\n", "\n")
    l = len(text.split(" "))
    if l > 15000 :
        return 'choose a smaller file'
    elif  l < 1000:
        return 'choose a larger file'

    format_i = '''
        Follow this format:

        [
            {"question": "here is the question",
            "answers": [
            "answer1", "answer2", "answer3", "answer4"
        ],
            "correct answer" : "only correct answer"
        }
        ,

        {"question": "here is the question",
            "answers": [
            "answer1", "answer2", "answer3", "answer4"
        ],
            "correct answer" : "only correct answer"
        }

        ]  
            ''' 
    prompt = f"create a {level} quiz with {num_qus} question and 4 answer and include the answer for the following text:\n\n{text}\n\n{format_i}"

    result = model.generate_content(
    prompt,
    generation_config={
        "max_output_tokens": 8192,
        "temperature": 0.2,
    },
  )
    token = len(text) + len(result.text)
    print(token)
    return eval(result.text)

if __name__ == '__main__':
    app.run(debug=False, port="8080")